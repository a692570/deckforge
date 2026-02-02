#!/usr/bin/env python3
"""
PPTX Export Script for Presentation Maker
Converts HTML presentations to PowerPoint (.pptx) format with full styling support.

Usage:
    python pptx_export.py input.html [--output output.pptx]

Features:
    - Parses HTML slide sections with CSS class detection
    - Maps slide types to appropriate master layouts
    - Extracts and applies CSS background colors
    - Embeds images with proper sizing and positioning
    - Typography: h1/h2/body with proper font sizes
    - Extracts speaker notes from HTML comments
    - 16:9 widescreen dimensions
    - Color theming from CSS variables
"""

import argparse
import base64
import io
import os
import re
import sys
from pathlib import Path
from urllib.parse import urlparse, unquote
from html.parser import HTMLParser

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
except ImportError:
    print("""
‚ùå ERROR: python-pptx is not installed!

To install, run one of these commands:
    pip3 install python-pptx
    pip install python-pptx

For more help: https://python-pptx.readthedocs.io/
""")
    sys.exit(1)


class SlideParser(HTMLParser):
    """Parse HTML to extract slide content, images, and metadata"""
    
    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_slide = None
        self.current_tag = None
        self.current_attrs = {}
        self.current_data = []
        self.in_speaker_notes = False
        self.in_style = False
        self.style_content = []
        self.css_vars = {}
        self.slide_counter = 0
        self.tag_stack = []
        
    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        self.tag_stack.append((tag, attrs_dict))
        
        # Capture style tag content for CSS parsing
        if tag == 'style':
            self.in_style = True
            return
        
        # Start of a new slide
        if tag == 'section' and 'slide' in attrs_dict.get('class', ''):
            self.slide_counter += 1
            classes = attrs_dict.get('class', '')
            
            # Determine slide type from CSS classes
            slide_type = 'content'
            layout = 'title_content'
            
            if 'title-slide' in classes:
                slide_type = 'title'
                layout = 'title_slide'
            elif 'split-slide' in classes or 'two-column' in classes:
                slide_type = 'split'
                layout = 'two_column'
            elif 'stats-slide' in classes or 'stats' in classes:
                slide_type = 'stats'
                layout = 'title_content'
            elif 'quote-slide' in classes or 'quote' in classes:
                slide_type = 'quote'
                layout = 'title_content'
            elif 'image-slide' in classes or 'image-heavy' in classes:
                slide_type = 'image'
                layout = 'title_content'
            elif 'section-header' in classes or 'section' in classes:
                slide_type = 'section'
                layout = 'section_header'
            elif 'closing-slide' in classes or 'closing' in classes:
                slide_type = 'closing'
                layout = 'title_slide'
            elif 'blank' in classes:
                slide_type = 'blank'
                layout = 'blank'
            
            self.current_slide = {
                'number': self.slide_counter,
                'type': slide_type,
                'layout': layout,
                'classes': classes,
                'title': '',
                'subtitle': '',
                'content': [],
                'images': [],
                'speaker_notes': '',
                'background': None
            }
        
        # Check for speaker notes div
        if tag == 'div' and 'speaker-notes' in attrs_dict.get('class', ''):
            self.in_speaker_notes = True
            self.current_data = []
            return
        
        # Track current tag for content extraction
        if self.current_slide and not self.in_speaker_notes:
            if tag in ['h1', 'h2', 'h3', 'p', 'li', 'strong', 'em', 'span']:
                self.current_tag = tag
                self.current_attrs = attrs_dict
                self.current_data = []
            
            # Extract image src
            if tag == 'img':
                src = attrs_dict.get('src', '')
                alt = attrs_dict.get('alt', '')
                if src:
                    self.current_slide['images'].append({
                        'src': src,
                        'alt': alt,
                        'style': attrs_dict.get('style', '')
                    })
    
    def handle_data(self, data):
        if self.in_style:
            self.style_content.append(data)
        elif self.current_tag or self.in_speaker_notes:
            self.current_data.append(data)
    
    def handle_endtag(self, tag):
        # End of style tag
        if tag == 'style':
            self.in_style = False
            self.css_vars = self._extract_css_variables(''.join(self.style_content))
            return
        
        if not self.current_slide:
            if self.tag_stack:
                self.tag_stack.pop()
            return
        
        content = ''.join(self.current_data).strip()
        
        # Speaker notes
        if tag == 'div' and self.in_speaker_notes:
            self.current_slide['speaker_notes'] = content
            self.in_speaker_notes = False
            self.current_data = []
            if self.tag_stack:
                self.tag_stack.pop()
            return
        
        # Title (H1)
        if tag == 'h1' and content:
            self.current_slide['title'] = content
        
        # Subtitle (H2) - could be subtitle or main title if no h1
        elif tag == 'h2' and content:
            if not self.current_slide['title']:
                self.current_slide['title'] = content
            else:
                self.current_slide['subtitle'] = content
        
        # H3 - often used as section headers within slides
        elif tag == 'h3' and content:
            self.current_slide['content'].append({
                'type': 'heading',
                'text': content,
                'level': 3
            })
        
        # Content paragraphs
        elif tag == 'p' and content:
            # Check if this is actually a subtitle (short text, no other content)
            if not self.current_slide['title'] and len(content) < 100:
                self.current_slide['title'] = content
            else:
                self.current_slide['content'].append({
                    'type': 'text',
                    'text': content
                })
        
        # List items
        elif tag == 'li' and content:
            self.current_slide['content'].append({
                'type': 'bullet',
                'text': content
            })
        
        # End of slide
        if tag == 'section' and self.current_slide:
            # Try to extract background from inline style or CSS
            self.current_slide['background'] = self._extract_background(self.current_slide['classes'])
            self.slides.append(self.current_slide)
            self.current_slide = None
        
        if tag == self.current_tag:
            self.current_tag = None
            self.current_attrs = {}
            self.current_data = []
        
        if self.tag_stack:
            self.tag_stack.pop()
    
    def _extract_css_variables(self, css_content):
        """Extract CSS variables from style content"""
        vars_dict = {}
        
        # Match :root CSS variables
        root_pattern = r':root\s*\{([^}]+)\}'
        root_match = re.search(root_pattern, css_content, re.DOTALL)
        
        if root_match:
            root_content = root_match.group(1)
            var_pattern = r'--([^:]+):\s*([^;]+);'
            for match in re.finditer(var_pattern, root_content):
                var_name = match.group(1).strip()
                var_value = match.group(2).strip()
                vars_dict[var_name] = var_value
        
        return vars_dict
    
    def _extract_background(self, classes):
        """Try to determine background color from common class patterns"""
        if 'title-slide' in classes or 'bg-dark' in classes or 'dark' in classes:
            return '#0A0A0A'
        if 'bg-light' in classes or 'light' in classes:
            return '#ffffff'
        return None


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    if not hex_color:
        return (255, 255, 255)
    
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except:
        return (255, 255, 255)


def rgb_color_from_hex(hex_color):
    """Create RGBColor from hex string"""
    rgb = hex_to_rgb(hex_color)
    return RGBColor(rgb[0], rgb[1], rgb[2])


def parse_color(color_value):
    """Parse various color formats to RGBColor"""
    if not color_value:
        return RGBColor(255, 255, 255)
    
    color_value = color_value.strip()
    
    # Hex color
    if color_value.startswith('#'):
        return rgb_color_from_hex(color_value)
    
    # RGB/RGBA
    rgb_match = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', color_value)
    if rgb_match:
        return RGBColor(
            int(rgb_match.group(1)),
            int(rgb_match.group(2)),
            int(rgb_match.group(3))
        )
    
    # Named colors (common ones)
    named_colors = {
        'white': '#ffffff',
        'black': '#000000',
        'red': '#ff0000',
        'green': '#00ff00',
        'blue': '#0000ff',
        'transparent': '#ffffff',
    }
    if color_value.lower() in named_colors:
        return rgb_color_from_hex(named_colors[color_value.lower()])
    
    return RGBColor(255, 255, 255)


def get_or_create_layout(prs, layout_name):
    """Get a slide layout by name, or return a suitable default"""
    layout_mapping = {
        'title_slide': ['Title Slide', 'Title Only', 'title'],
        'title_content': ['Title and Content', 'Title Content', 'content'],
        'section_header': ['Section Header', 'Section', 'section'],
        'two_column': ['Two Content', 'Two Column', 'Comparison', 'comparison'],
        'blank': ['Blank', 'blank'],
    }
    
    search_names = layout_mapping.get(layout_name, [layout_name])
    
    # Try to find matching layout
    for layout in prs.slide_layouts:
        layout_name_lower = layout.name.lower()
        for search in search_names:
            if search.lower() in layout_name_lower:
                return layout
    
    # Fallback to index-based selection
    if layout_name == 'title_slide' and len(prs.slide_layouts) > 0:
        return prs.slide_layouts[0]
    elif layout_name == 'title_content' and len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    elif layout_name == 'section_header' and len(prs.slide_layouts) > 2:
        return prs.slide_layouts[2]
    elif layout_name == 'two_column' and len(prs.slide_layouts) > 3:
        return prs.slide_layouts[3]
    else:
        return prs.slide_layouts[-1] if prs.slide_layouts else prs.slide_layouts[0]


def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    if not color:
        return
    
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = parse_color(color)
    except Exception as e:
        print(f"Warning: Could not set background color: {e}")


def add_image_to_slide(slide, img_path, left=None, top=None, width=None, height=None):
    """Add an image to a slide with optional positioning"""
    try:
        if left is None:
            left = Inches(8)
        if top is None:
            top = Inches(1.5)
        if height is None and width is None:
            height = Inches(4)
        
        return slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    except Exception as e:
        print(f"Warning: Could not add image {img_path}: {e}")
        return None


def resolve_image_path(src, html_dir):
    """Resolve image source to absolute path"""
    if not src:
        return None
    
    # Data URI - decode and save to temp file
    if src.startswith('data:'):
        try:
            # Parse data URI
            match = re.match(r'data:image/(\w+);base64,(.+)', src)
            if match:
                ext = match.group(1)
                data = match.group(2)
                img_data = base64.b64decode(data)
                
                # Create temp file
                temp_path = os.path.join(html_dir, f'_temp_image.{ext}')
                with open(temp_path, 'wb') as f:
                    f.write(img_data)
                return temp_path
        except Exception as e:
            print(f"Warning: Could not decode data URI: {e}")
            return None
    
    # Remote URL - skip (would need download logic)
    if src.startswith(('http://', 'https://')):
        print(f"Note: Remote image skipped (download not implemented): {src[:50]}...")
        return None
    
    # Relative path
    if src.startswith('/'):
        # Absolute from root - try relative to html_dir
        src = src.lstrip('/')
    
    img_path = os.path.join(html_dir, src)
    if os.path.exists(img_path):
        return img_path
    
    # Try common variations
    variations = [
        src,
        src.lstrip('./'),
        os.path.basename(src),
        os.path.join('assets', os.path.basename(src)),
        os.path.join('images', os.path.basename(src)),
    ]
    
    for var in variations:
        test_path = os.path.join(html_dir, var)
        if os.path.exists(test_path):
            return test_path
    
    return None


def apply_text_formatting(text_frame, color, font_size=18, bold=False, font_name='Arial'):
    """Apply consistent formatting to a text frame"""
    try:
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name
            paragraph.font.color.rgb = color
            paragraph.font.bold = bold
            
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.color.rgb = color
                run.font.bold = bold
    except Exception as e:
        print(f"Warning: Could not apply text formatting: {e}")


def create_presentation(slides, css_vars, output_path, html_dir, verbose=False):
    """Create a PowerPoint presentation from parsed slides"""
    
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Extract colors from CSS variables with fallbacks
    bg_primary = css_vars.get('bg-primary', '#ffffff')
    bg_dark = css_vars.get('bg-dark', '#0A0A0A')
    text_primary = css_vars.get('text-primary', '#1a1a1a')
    text_on_dark = css_vars.get('text-on-dark', '#ffffff')
    accent = css_vars.get('accent', '#00E3AA')
    
    # Font preferences from CSS
    font_display = css_vars.get('font-display', 'Arial').replace("'", "").replace('"', '')
    font_body = css_vars.get('font-body', 'Arial').replace("'", "").replace('"', '')
    
    temp_files = []  # Track temp files for cleanup
    
    for i, slide_data in enumerate(slides):
        if verbose:
            print(f"  Processing slide {i+1}/{len(slides)}: {slide_data['type']} - {slide_data['title'][:40]}...")
        
        # Determine colors based on slide type/background
        if slide_data['type'] in ['title', 'section', 'closing'] or 'dark' in slide_data['classes']:
            bg_color = bg_dark
            text_color = text_on_dark
        else:
            bg_color = bg_primary
            text_color = text_primary
        
        # Override if specific background was detected
        if slide_data['background']:
            bg_color = slide_data['background']
            text_color = text_on_dark if bg_color in ['#0A0A0A', '#000000', '#0f172a', '#030712'] else text_primary
        
        # Get appropriate layout
        layout = get_or_create_layout(prs, slide_data['layout'])
        
        # Add slide
        slide = prs.slides.add_slide(layout)
        
        # Set background color
        set_slide_background(slide, bg_color)
        
        # Convert colors
        text_rgb = parse_color(text_color)
        accent_rgb = parse_color(accent)
        
        # Add title
        if slide.shapes.title and slide_data['title']:
            title = slide.shapes.title
            title.text = slide_data['title']
            
            # Title slide gets larger font
            title_size = 54 if slide_data['type'] == 'title' else 40
            apply_text_formatting(title.text_frame, text_rgb, title_size, bold=True, font_name=font_display)
            
            # Center align title on title slides
            if slide_data['type'] == 'title':
                for paragraph in title.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
        
        # Add subtitle (title slides)
        if slide_data['type'] in ['title', 'closing'] and slide_data['subtitle']:
            # Try to find subtitle placeholder
            for shape in slide.placeholders:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # SUBTITLE
                    shape.text = slide_data['subtitle']
                    apply_text_formatting(shape.text_frame, text_rgb, 24, font_name=font_body)
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                    break
        
        # Add content
        if slide_data['content']:
            # Find the content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    # BODY, OBJECT, VERTICAL_BODY, VERTICAL_OBJECT
                    if ph_type in [1, 12, 13, 14, 15]:
                        content_placeholder = shape
                        break
            
            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.clear()
                
                for j, item in enumerate(slide_data['content']):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = item['text']
                    p.font.color.rgb = text_rgb
                    p.font.size = Pt(20) if item['type'] == 'heading' else Pt(18)
                    p.font.name = font_display if item['type'] == 'heading' else font_body
                    p.font.bold = item['type'] == 'heading'
                    
                    if item['type'] == 'bullet':
                        p.level = 0
                    elif item['type'] == 'heading':
                        p.level = 0
                        p.space_after = Pt(6)
                    else:
                        p.level = 0
            else:
                # No placeholder found, add text box manually
                if slide_data['content']:
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(11)
                    height = Inches(4)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    tf.word_wrap = True
                    
                    for j, item in enumerate(slide_data['content']):
                        if j == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        
                        p.text = item['text']
                        p.font.color.rgb = text_rgb
                        p.font.size = Pt(20) if item['type'] == 'heading' else Pt(18)
                        p.font.name = font_display if item['type'] == 'heading' else font_body
                        p.font.bold = item['type'] == 'heading'
        
        # Add images
        if slide_data['images']:
            for img_data in slide_data['images']:
                img_path = resolve_image_path(img_data['src'], html_dir)
                
                if img_path:
                    try:
                        # Determine positioning based on slide type
                        if slide_data['type'] == 'image' or slide_data['layout'] == 'two_column':
                            # Image takes right half
                            left = Inches(6.5)
                            top = Inches(1)
                            height = Inches(5.5)
                        elif slide_data['type'] == 'title':
                            # Logo/accent image, smaller
                            left = Inches(5.5)
                            top = Inches(5.5)
                            height = Inches(1.5)
                        else:
                            # Default right-side placement
                            left = Inches(8)
                            top = Inches(1.5)
                            height = Inches(4)
                        
                        pic = add_image_to_slide(slide, img_path, left=left, top=top, height=height)
                        
                        if pic and img_path.startswith(os.path.join(html_dir, '_temp')):
                            temp_files.append(img_path)
                            
                    except Exception as e:
                        print(f"Warning: Could not add image {img_data['src']}: {e}")
        
        # Add speaker notes
        if slide_data['speaker_notes']:
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data['speaker_notes']
            except Exception as e:
                if verbose:
                    print(f"Warning: Could not add speaker notes: {e}")
    
    # Clean up temp files
    for temp_file in temp_files:
        try:
            os.remove(temp_file)
        except:
            pass
    
    # Save presentation
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description='Convert HTML presentations to PowerPoint (.pptx) format with full styling support',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
    python pptx_export.py presentation.html
    python pptx_export.py presentation.html -o output.pptx -v
        """
    )
    parser.add_argument('input', help='Input HTML file path')
    parser.add_argument('--output', '-o', help='Output PPTX file path (optional)')
    parser.add_argument('--verbose', '-v', action='store_true', help='Verbose output')
    
    args = parser.parse_args()
    
    # Validate input
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"‚ùå Error: Input file not found: {args.input}")
        sys.exit(1)
    
    if not input_path.suffix.lower() in ['.html', '.htm']:
        print(f"‚ö†Ô∏è  Warning: Input file doesn't have .html extension: {args.input}")
    
    # Determine output path
    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_suffix('.pptx')
    
    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    print(f"üìÑ Input: {input_path.absolute()}")
    print(f"üìä Output: {output_path.absolute()}")
    
    # Read HTML file
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
    except Exception as e:
        print(f"‚ùå Error reading input file: {e}")
        sys.exit(1)
    
    # Parse slides
    parser = SlideParser()
    parser.feed(html_content)
    slides = parser.slides
    css_vars = parser.css_vars
    
    if not slides:
        print("‚ùå Error: No slides found in HTML file")
        print("   Make sure your HTML has <section class='slide'> elements")
        sys.exit(1)
    
    print(f"üé® Found {len(css_vars)} CSS variables")
    print(f"üìë Found {len(slides)} slides")
    
    if args.verbose:
        for slide in slides:
            print(f"   - Slide {slide['number']}: {slide['type']} | {slide['title'][:30]}... | {len(slide['content'])} content items | {len(slide['images'])} images")
    
    # Create presentation
    try:
        html_dir = str(input_path.parent)
        create_presentation(slides, css_vars, str(output_path), html_dir, args.verbose)
        
        file_size = output_path.stat().st_size
        print(f"‚úÖ Successfully created: {output_path}")
        print(f"   Slides: {len(slides)}")
        print(f"   File size: {file_size / 1024:.1f} KB")
        
    except Exception as e:
        print(f"‚ùå Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
